"""Generate PowerPoint slides from set_theoretic_analysis.md."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette — Light theme
BG_DARK = RGBColor(0xFF, 0xFF, 0xFF)      # white background
BG_MID = RGBColor(0xF2, 0xF2, 0xF7)      # light gray cards
ACCENT = RGBColor(0x00, 0x7A, 0x6D)      # dark teal
ACCENT2 = RGBColor(0x26, 0x5C, 0x9E)     # dark blue
ACCENT3 = RGBColor(0x4A, 0x4A, 0x2A)     # dark olive (code text)
WHITE = RGBColor(0x2D, 0x2D, 0x2D)       # dark text (replaces white)
GRAY = RGBColor(0x66, 0x66, 0x66)        # medium gray text
RED_SOFT = RGBColor(0xC0, 0x39, 0x2B)    # red
GREEN_SOFT = RGBColor(0x1E, 0x7D, 0x34)  # green

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Consolas"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet(tf, text, level=0, font_size=16, color=WHITE, bold=False, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_after = Pt(4)
    return p

def add_code_box(slide, left, top, width, height, code, title="", title_color=ACCENT):
    from pptx.util import Emu
    # Background box
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xFA)
    shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shape.line.width = Pt(0.75)

    y_offset = top + 0.15
    if title:
        add_text_box(slide, left + 0.2, y_offset, width - 0.4, 0.35,
                     title, font_size=13, color=title_color, bold=True, font_name="Calibri")
        y_offset += 0.3

    add_text_box(slide, left + 0.2, y_offset, width - 0.4, height - (y_offset - top) - 0.1,
                 code, font_size=12, color=ACCENT3, font_name="Consolas")


# ========================================================================
# SLIDE 1: Title + Core Design Decision
# ========================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, BG_DARK)

# Title
add_text_box(slide1, 0.6, 0.3, 12, 0.8,
             "Set-Theoretic Abstraction in Bitmap Verification",
             font_size=36, color=ACCENT, bold=True, font_name="Calibri")

# Subtitle
add_text_box(slide1, 0.6, 1.1, 12, 0.5,
             "Why treating BitmapView as a first-class mathematical object matters",
             font_size=20, color=GRAY, font_name="Calibri")

# Divider
shape = slide1.shapes.add_shape(1, Inches(0.6), Inches(1.7), Inches(12), Inches(0.02))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Left column: bitmap_raw
add_text_box(slide1, 0.6, 1.9, 5.5, 0.4,
             "bitmap_raw  — Passive View", font_size=22, color=RED_SOFT, bold=True, font_name="Calibri")

tf1 = add_text_box(slide1, 0.6, 2.4, 5.5, 1.2,
                   "BitmapView is just the return type of view()",
                   font_size=16, color=GRAY, font_name="Calibri")
add_bullet(tf1, "All predicates & lemmas live on Bitmap (concrete)", font_size=15, color=WHITE)
add_bullet(tf1, "Abstract lemmas require inv() (full concrete invariant)", font_size=15, color=WHITE)
add_bullet(tf1, "is_bit_set defined on BOTH types → bridging lemma needed", font_size=15, color=WHITE)
add_bullet(tf1, "~1169 lines of proofs, all on impl Bitmap", font_size=15, color=GRAY)

add_code_box(slide1, 0.6, 4.2, 5.5, 2.0,
             '// Predicate bound to Bitmap\n'
             'impl Bitmap {\n'
             '  pub open spec fn has_free_range_at(\n'
             '    &self, start: int, n: int\n'
             '  ) -> bool {\n'
             '    self.all_bits_unset_in_range(start, start+n)\n'
             '  }\n'
             '  // Lemma requires inv()\n'
             '  proof fn lemma_usage_bound(&self, p, n)\n'
             '    requires self.inv(), ...\n'
             '}',
             title="bitmap_raw: specs on concrete Bitmap", title_color=RED_SOFT)

# Right column: bitmap_new
add_text_box(slide1, 7.0, 1.9, 5.8, 0.4,
             "bitmap_new  — Active Specification Model", font_size=22, color=GREEN_SOFT, bold=True, font_name="Calibri")

tf2 = add_text_box(slide1, 7.0, 2.4, 5.8, 1.2,
                   "BitmapView is a self-contained mathematical object",
                   font_size=16, color=GRAY, font_name="Calibri")
add_bullet(tf2, "Predicates & abstract lemmas live on BitmapView", font_size=15, color=WHITE)
add_bullet(tf2, "Abstract lemmas require only wf() (pure set constraint)", font_size=15, color=WHITE)
add_bullet(tf2, "is_bit_set defined ONCE on BitmapView — no bridging", font_size=15, color=WHITE)
add_bullet(tf2, "~716 lines of proofs on Bitmap + ~394 on BitmapView", font_size=15, color=GRAY)

add_code_box(slide1, 7.0, 4.2, 5.8, 2.0,
             '// Predicate on abstract model\n'
             'impl BitmapView {\n'
             '  pub open spec fn has_free_range_at(\n'
             '    &self, start: int, n: int\n'
             '  ) -> bool {\n'
             '    self.all_bits_unset_in_range(start, start+n)\n'
             '  }\n'
             '  // Lemma requires only wf()\n'
             '  proof fn lemma_usage_bound(&self, p, n)\n'
             '    requires self.wf(), ...\n'
             '}',
             title="bitmap_new: specs on abstract BitmapView", title_color=GREEN_SOFT)

# Bottom tagline
add_text_box(slide1, 0.6, 6.5, 12, 0.5,
             '"The view is a projection of the struct"  vs  "The struct is a realization of the model"',
             font_size=17, color=ACCENT, bold=False, font_name="Calibri", alignment=PP_ALIGN.CENTER)


# ========================================================================
# SLIDE 2: Key Advantages (3-column layout)
# ========================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)

add_text_box(slide2, 0.6, 0.3, 12, 0.7,
             "Key Advantages of Set-Theoretic Abstraction",
             font_size=32, color=ACCENT, bold=True, font_name="Calibri")

shape = slide2.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(12), Inches(0.02))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Column 1: Representation Independence
col1_x = 0.6
add_text_box(slide2, col1_x, 1.2, 3.8, 0.35,
             "① Representation Independence", font_size=18, color=ACCENT2, bold=True, font_name="Calibri")
tf = add_text_box(slide2, col1_x, 1.6, 3.8, 2.2,
                  "", font_size=14, color=WHITE, font_name="Calibri")
add_bullet(tf, "Swap RawArray<u8> → Vec<u64>?", font_size=14, color=WHITE, bold=True)
add_bullet(tf, "15+ abstract lemmas survive unchanged", level=1, font_size=13, color=GREEN_SOFT)
add_bullet(tf, "Only re-prove byte↔set coupling", level=1, font_size=13, color=GRAY)
add_bullet(tf, "", font_size=6)
add_bullet(tf, "In bitmap_raw: must re-prove ALL", font_size=14, color=RED_SOFT)
add_bullet(tf, "~40 lemmas tied to impl Bitmap", level=1, font_size=13, color=GRAY)

# Column 2: Weaker Preconditions
col2_x = 4.8
add_text_box(slide2, col2_x, 1.2, 3.8, 0.35,
             "② Weaker Preconditions", font_size=18, color=ACCENT2, bold=True, font_name="Calibri")
tf = add_text_box(slide2, col2_x, 1.6, 3.8, 2.2,
                  "", font_size=14, color=WHITE, font_name="Calibri")
add_bullet(tf, "Lemmas need wf(), not inv()", font_size=14, color=WHITE, bold=True)
add_bullet(tf, "wf() = ∀i ∈ set_bits → 0 ≤ i < n", level=1, font_size=13, color=GREEN_SOFT)
add_bullet(tf, "Callable in intermediate states", level=1, font_size=13, color=GREEN_SOFT)
add_bullet(tf, "", font_size=6)
add_bullet(tf, "In bitmap_raw: need full inv()", font_size=14, color=RED_SOFT)
add_bullet(tf, "Includes bytes, usage, next_free...", level=1, font_size=13, color=GRAY)
add_bullet(tf, "Can't call during bit modification", level=1, font_size=13, color=GRAY)

# Column 3: Eliminated Duplication
col3_x = 9.0
add_text_box(slide2, col3_x, 1.2, 3.8, 0.35,
             "③ No Impedance Mismatch", font_size=18, color=ACCENT2, bold=True, font_name="Calibri")
tf = add_text_box(slide2, col3_x, 1.6, 3.8, 2.2,
                  "", font_size=14, color=WHITE, font_name="Calibri")
add_bullet(tf, "is_bit_set: ONE definition", font_size=14, color=WHITE, bold=True)
add_bullet(tf, "= set_bits.contains(index)", level=1, font_size=13, color=GREEN_SOFT)
add_bullet(tf, "No bridging lemma needed", level=1, font_size=13, color=GREEN_SOFT)
add_bullet(tf, "", font_size=6)
add_bullet(tf, "In bitmap_raw: TWO definitions", font_size=14, color=RED_SOFT)
add_bullet(tf, "Bitmap version bundles bounds check", level=1, font_size=13, color=GRAY)
add_bullet(tf, "Needs lemma_is_bit_set_equals_view", level=1, font_size=13, color=GRAY)

# Lower section: code comparison
add_code_box(slide2, 0.6, 4.2, 3.8, 1.8,
             '// bitmap_raw: call on Bitmap\n'
             'self.lemma_free_range_implies\n'
             '  _usage_bound(p, size);\n\n'
             '// Requires: self.inv() ← full\n'
             '//   concrete invariant',
             title="Lemma call site (raw)", title_color=RED_SOFT)

add_code_box(slide2, 4.8, 4.2, 3.8, 1.8,
             '// bitmap_new: call on View\n'
             'self@.lemma_free_range_implies\n'
             '  _usage_bound(p, size);\n\n'
             '// Requires: self@.wf() ← pure\n'
             '//   set-theoretic constraint',
             title="Lemma call site (new)", title_color=GREEN_SOFT)

add_code_box(slide2, 9.0, 4.2, 3.8, 1.8,
             '// Open invariant → auto wf()\n'
             'pub open spec fn inv(&self)\n'
             '  -> bool {\n'
             '  &&& self@.wf()    // visible!\n'
             '  &&& self.internal_inv()\n'
             '}',
             title="Invariant structure (new)", title_color=ACCENT)

# Bottom row: additional advantages
add_text_box(slide2, 0.6, 6.3, 12, 0.35,
             "Also:  ④ Composable set algebra (vstd directly usable)   "
             "⑤ Physical proof stratification (spec vs proof files)   "
             "⑥ SMT solver extracts wf() automatically from open inv()",
             font_size=14, color=GRAY, font_name="Calibri", alignment=PP_ALIGN.CENTER)


# ========================================================================
# SLIDE 3: Architecture & Summary
# ========================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)

add_text_box(slide3, 0.6, 0.3, 12, 0.7,
             "Proof Architecture & Impact",
             font_size=32, color=ACCENT, bold=True, font_name="Calibri")

shape = slide3.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(12), Inches(0.02))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Left: Layer diagram
add_text_box(slide3, 0.6, 1.2, 5.5, 0.35,
             "Proof Stratification", font_size=22, color=ACCENT2, bold=True, font_name="Calibri")

layers = [
    ("Layer 3: Client Code", "Interacts with BitmapView predicates + inv()", GRAY),
    ("Layer 2: Representation Coupling", "byte ops ↔ set_bits changes  (lib.proof.rs)", ACCENT2),
    ("Layer 1: Abstract Specification", "set-theoretic properties  (lib.spec.rs)", GREEN_SOFT),
    ("Layer 0: vstd", "set_lib, set axioms", GRAY),
]

y = 1.7
for title, desc, color in layers:
    box_shape = slide3.shapes.add_shape(
        1, Inches(0.8), Inches(y), Inches(5.0), Inches(0.7))
    box_shape.fill.solid()
    box_shape.fill.fore_color.rgb = BG_MID
    box_shape.line.color.rgb = color
    box_shape.line.width = Pt(1.5)

    tf = box_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY
    p2.font.name = "Calibri"

    # Arrow between layers
    if y > 1.7:
        arrow_y = y - 0.15
        add_text_box(slide3, 3.0, arrow_y, 0.5, 0.2, "▲", font_size=12, color=GRAY,
                     alignment=PP_ALIGN.CENTER, font_name="Calibri")

    y += 0.85

# Right: Summary table
add_text_box(slide3, 7.0, 1.2, 5.8, 0.35,
             "Impact Summary", font_size=22, color=ACCENT2, bold=True, font_name="Calibri")

# Table
from pptx.util import Inches, Pt
table_shape = slide3.shapes.add_table(8, 3, Inches(7.0), Inches(1.7), Inches(5.8), Inches(3.5))
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(2.6)
table.columns[1].width = Inches(1.4)
table.columns[2].width = Inches(1.8)

headers = ["Property", "raw", "new"]
rows = [
    ["Repr. change re-proves abstracts?", "All", "None"],
    ["Abstract lemma precondition", "inv()", "wf()"],
    ["is_bit_set definitions", "2 + bridge", "1"],
    ["Hypothetical BitmapView in proofs?", "No", "Yes"],
    ["vstd set_lib usage", "Wrapped", "Direct"],
    ["Proof boundary", "Mixed", "Separated"],
    ["Proof lines (lib.proof.rs)", "1169", "716 (−39%)"],
]

for col_idx, h in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.font.name = "Calibri"
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xF0)

for row_idx, row in enumerate(rows):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.name = "Calibri"
            if col_idx == 1:
                p.font.color.rgb = RED_SOFT
            elif col_idx == 2:
                p.font.color.rgb = GREEN_SOFT
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xFC) if row_idx % 2 == 0 else RGBColor(0xEE, 0xEE, 0xF5)

# Bottom tagline
add_text_box(slide3, 0.6, 6.2, 12, 0.8,
             "Bottom Line:  Separating what a bitmap means (a bounded set of integers)\n"
             "from how it's implemented (a byte array with bookkeeping)\n"
             "yields reusable proofs, better solver performance, and cleaner maintenance.",
             font_size=16, color=ACCENT, font_name="Calibri", alignment=PP_ALIGN.CENTER)

# ========================================================================
# SLIDE 4: Outline — Improving Specs via Test Generation
# ========================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_DARK)

add_text_box(slide4, 0.6, 0.3, 12, 0.7,
             "Improving Specs via Test Generation — Outline",
             font_size=32, color=ACCENT, bold=True, font_name="Calibri")

shape = slide4.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(12), Inches(0.02))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Goal
add_text_box(slide4, 0.8, 1.2, 11.5, 0.5,
             "Goal:  Improve the correctness and completeness of specifications via test generation",
             font_size=22, color=ACCENT2, bold=True, font_name="Calibri")

add_text_box(slide4, 0.8, 1.75, 11.5, 0.4,
             "Plans for VeruSage and Nanvix",
             font_size=18, color=GRAY, font_name="Calibri")

# Steps as cards
steps = [
    ("Step 1", "Build test execution environment",
     "Set up infrastructure for running generated tests against Verus specs",
     ACCENT2),
    ("Step 2", "Generate docstrings & documents from implementation",
     "Extract structured documentation from existing Nanvix code to guide LLM generation",
     ACCENT2),
    ("Step 3", "Generate specs and tests; use tests to evaluate specs",
     "LLMs produce candidate specifications + test suites; tests validate spec correctness",
     ACCENT2),
]

y = 2.3
for label, title, desc, color in steps:
    # Card background
    card = slide4.shapes.add_shape(
        1, Inches(0.8), Inches(y), Inches(7.5), Inches(0.9))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_MID
    card.line.color.rgb = color
    card.line.width = Pt(1.5)

    # Step label
    add_text_box(slide4, 1.0, y + 0.08, 1.0, 0.35,
                 label, font_size=14, color=RGBColor(0xFF, 0xFF, 0xFF),
                 bold=True, font_name="Calibri")
    # Step label background pill
    pill = slide4.shapes.add_shape(
        1, Inches(0.9), Inches(y + 0.1), Inches(0.9), Inches(0.3))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    tf_pill = pill.text_frame
    tf_pill.word_wrap = False
    p_pill = tf_pill.paragraphs[0]
    p_pill.text = label
    p_pill.font.size = Pt(13)
    p_pill.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p_pill.font.bold = True
    p_pill.font.name = "Calibri"
    p_pill.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide4, 1.9, y + 0.05, 6.2, 0.35,
                 title, font_size=17, color=WHITE, bold=True, font_name="Calibri")
    # Description
    add_text_box(slide4, 1.9, y + 0.45, 6.2, 0.35,
                 desc, font_size=13, color=GRAY, font_name="Calibri")

    y += 1.05

# Optional items on the right
add_text_box(slide4, 8.8, 2.3, 4.0, 0.35,
             "Optional / Future Work", font_size=18, color=ACCENT, bold=True, font_name="Calibri")

opt_card = slide4.shapes.add_shape(
    1, Inches(8.8), Inches(2.8), Inches(4.0), Inches(1.2))
opt_card.fill.solid()
opt_card.fill.fore_color.rgb = BG_MID
opt_card.line.color.rgb = ACCENT
opt_card.line.width = Pt(1)

tf_opt = add_text_box(slide4, 9.0, 2.9, 3.6, 0.3,
                      "Rust library specification inference for Verus",
                      font_size=14, color=WHITE, bold=False, font_name="Calibri")
add_bullet(tf_opt, "Infer specs for std/external crates", font_size=13, color=GRAY)

opt_card2 = slide4.shapes.add_shape(
    1, Inches(8.8), Inches(4.2), Inches(4.0), Inches(1.2))
opt_card2.fill.solid()
opt_card2.fill.fore_color.rgb = BG_MID
opt_card2.line.color.rgb = ACCENT
opt_card2.line.width = Pt(1)

tf_opt2 = add_text_box(slide4, 9.0, 4.3, 3.6, 0.3,
                       "Learn the abstract level during spec generation",
                       font_size=14, color=WHITE, bold=False, font_name="Calibri")
add_bullet(tf_opt2, "Train LLMs to choose the right abstraction layer",
           font_size=13, color=GRAY)
add_bullet(tf_opt2, "e.g., BitmapView-level vs Bitmap-level specs",
           font_size=13, color=GRAY)


# ========================================================================
# SLIDE 5: Why LLM-generated tests can improve specs?
# ========================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_DARK)

add_text_box(slide5, 0.6, 0.3, 12, 0.7,
             "Why LLM-Generated Tests Can Improve Specs",
             font_size=32, color=ACCENT, bold=True, font_name="Calibri")

shape = slide5.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(12), Inches(0.02))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# Core idea
add_text_box(slide5, 0.8, 1.2, 11.5, 0.5,
             "Core Idea:  Leverage different LLMs to generate specs and tests independently",
             font_size=22, color=ACCENT2, bold=True, font_name="Calibri")

# Approach 1 — Left card
card1 = slide5.shapes.add_shape(
    1, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.8))
card1.fill.solid()
card1.fill.fore_color.rgb = BG_MID
card1.line.color.rgb = RED_SOFT
card1.line.width = Pt(2)

add_text_box(slide5, 1.0, 2.15, 5.0, 0.4,
             "① Inconsistency Detection  (DeepTest)",
             font_size=20, color=RED_SOFT, bold=True, font_name="Calibri")

tf1 = add_text_box(slide5, 1.0, 2.65, 5.0, 0.3,
                   "", font_size=15, color=WHITE, font_name="Calibri")
add_bullet(tf1, "LLM-A generates a specification", font_size=15, color=WHITE)
add_bullet(tf1, "LLM-B generates tests for the same function", font_size=15, color=WHITE)
add_bullet(tf1, "", font_size=6)
add_bullet(tf1, "If spec and tests are inconsistent:", font_size=15, color=WHITE, bold=True)
add_bullet(tf1, "→ At least one of them is incorrect", level=1, font_size=14, color=RED_SOFT)
add_bullet(tf1, "→ Pinpoints where the spec needs fixing", level=1, font_size=14, color=RED_SOFT)
add_bullet(tf1, "", font_size=6)
add_bullet(tf1, "Cross-validation without ground truth!", font_size=15, color=ACCENT, bold=True)

# Diagram in card 1
add_text_box(slide5, 1.2, 4.6, 5.0, 0.6,
             "  LLM-A → Spec ─┐\n"
             "                 ├─ Compare → Find bugs\n"
             "  LLM-B → Tests ┘",
             font_size=13, color=ACCENT2, font_name="Consolas")

# Approach 2 — Right card
card2 = slide5.shapes.add_shape(
    1, Inches(6.8), Inches(2.0), Inches(5.8), Inches(3.8))
card2.fill.solid()
card2.fill.fore_color.rgb = BG_MID
card2.line.color.rgb = GREEN_SOFT
card2.line.width = Pt(2)

add_text_box(slide5, 7.0, 2.15, 5.4, 0.4,
             "② Majority Voting + NL Round-Trip",
             font_size=20, color=GREEN_SOFT, bold=True, font_name="Calibri")

tf2 = add_text_box(slide5, 7.0, 2.65, 5.4, 0.3,
                   "", font_size=15, color=WHITE, font_name="Calibri")
add_bullet(tf2, "Multiple LLMs generate test suites", font_size=15, color=WHITE)
add_bullet(tf2, "Majority voting: consensus = likely correct", font_size=15, color=WHITE)
add_bullet(tf2, "", font_size=6)
add_bullet(tf2, "NL round-trip validation:", font_size=15, color=WHITE, bold=True)
add_bullet(tf2, "1. Start with NL description of behavior",
           level=1, font_size=14, color=GREEN_SOFT)
add_bullet(tf2, "2. LLM generates test code", level=1, font_size=14, color=GREEN_SOFT)
add_bullet(tf2, "3. Translate tests back to NL", level=1, font_size=14, color=GREEN_SOFT)
add_bullet(tf2, "4. Compare to original NL description", level=1, font_size=14, color=GREEN_SOFT)
add_bullet(tf2, "", font_size=6)
add_bullet(tf2, "Semantic drift detection!", font_size=15, color=ACCENT, bold=True)

# Diagram in card 2
add_text_box(slide5, 7.2, 4.7, 5.2, 0.6,
             "  NL desc → LLM → Tests → Back to NL\n"
             "    ↑                          ↓\n"
             "    └──── Compare for drift ───┘",
             font_size=13, color=ACCENT2, font_name="Consolas")

# Bottom
add_text_box(slide5, 0.6, 6.2, 12, 0.5,
             "Tests serve as an independent oracle — if they disagree with specs,\n"
             "we know where to look without manual inspection.",
             font_size=16, color=ACCENT, font_name="Calibri", alignment=PP_ALIGN.CENTER)


# ========================================================================
# SLIDE 6: Challenges & Plans — Learning Abstract Level
# ========================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, BG_DARK)

add_text_box(slide6, 0.6, 0.3, 12, 0.7,
             "Learning the Abstract Level During Spec Generation",
             font_size=32, color=ACCENT, bold=True, font_name="Calibri")

shape = slide6.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(12), Inches(0.02))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# --- Challenge 1 (left) ---
ch1_card = slide6.shapes.add_shape(
    1, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.3))
ch1_card.fill.solid()
ch1_card.fill.fore_color.rgb = BG_MID
ch1_card.line.color.rgb = RED_SOFT
ch1_card.line.width = Pt(2)

add_text_box(slide6, 1.0, 1.4, 5.0, 0.4,
             "Challenge 1", font_size=20, color=RED_SOFT, bold=True, font_name="Calibri")
add_text_box(slide6, 1.0, 1.85, 5.0, 0.4,
             "The right abstract level is highly domain-specific",
             font_size=17, color=WHITE, bold=False, font_name="Calibri")

tf_ch1 = add_text_box(slide6, 1.0, 2.3, 5.0, 0.3,
                      "", font_size=14, color=WHITE, font_name="Calibri")
add_bullet(tf_ch1, "Bitmap: Set<int> vs raw byte array", font_size=14, color=GRAY)
add_bullet(tf_ch1, "File system: inode tree vs block layout", font_size=14, color=GRAY)
add_bullet(tf_ch1, "No universal rule — depends on verification goal", font_size=14, color=GRAY)

# --- Plan 1 (left, below) ---
pl1_card = slide6.shapes.add_shape(
    1, Inches(0.8), Inches(3.8), Inches(5.5), Inches(2.8))
pl1_card.fill.solid()
pl1_card.fill.fore_color.rgb = BG_MID
pl1_card.line.color.rgb = GREEN_SOFT
pl1_card.line.width = Pt(2)

add_text_box(slide6, 1.0, 3.9, 5.0, 0.4,
             "Plan 1:  Bottom-Up Abstraction with LLMs",
             font_size=18, color=GREEN_SOFT, bold=True, font_name="Calibri")

tf_pl1 = add_text_box(slide6, 1.0, 4.35, 5.0, 0.3,
                      "", font_size=14, color=WHITE, font_name="Calibri")
add_bullet(tf_pl1, "Start: generate low-level spec (bitmap_raw style)",
           font_size=14, color=WHITE)
add_bullet(tf_pl1, "Iterate: LLM proposes abstraction lifts", font_size=14, color=WHITE)
add_bullet(tf_pl1, "Move predicates from Bitmap → BitmapView",
           level=1, font_size=13, color=GREEN_SOFT)
add_bullet(tf_pl1, "Split inv() into open wf() + closed internal_inv()",
           level=1, font_size=13, color=GREEN_SOFT)
add_bullet(tf_pl1, "Verify: Re-check proofs after each lift",
           font_size=14, color=WHITE)

# Arrow/flow
add_text_box(slide6, 1.5, 6.0, 4.0, 0.4,
             "bitmap_raw  →  LLM refactor  →  bitmap_new",
             font_size=14, color=ACCENT2, bold=True, font_name="Consolas",
             alignment=PP_ALIGN.CENTER)

# --- Challenge 2 (right) ---
ch2_card = slide6.shapes.add_shape(
    1, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.3))
ch2_card.fill.solid()
ch2_card.fill.fore_color.rgb = BG_MID
ch2_card.line.color.rgb = RED_SOFT
ch2_card.line.width = Pt(2)

add_text_box(slide6, 7.0, 1.4, 5.4, 0.4,
             "Challenge 2", font_size=20, color=RED_SOFT, bold=True, font_name="Calibri")
add_text_box(slide6, 7.0, 1.85, 5.4, 0.4,
             "Abstraction quality is difficult to evaluate",
             font_size=17, color=WHITE, bold=False, font_name="Calibri")

tf_ch2 = add_text_box(slide6, 7.0, 2.3, 5.4, 0.3,
                      "", font_size=14, color=WHITE, font_name="Calibri")
add_bullet(tf_ch2, "\"Is this the right abstraction?\" — no automated metric",
           font_size=14, color=GRAY)
add_bullet(tf_ch2, "Proof size? Solver time? Reusability? All partial proxies",
           font_size=14, color=GRAY)
add_bullet(tf_ch2, "Human expert judgment is expensive and subjective",
           font_size=14, color=GRAY)

# --- Plan 2 (right, below) ---
pl2_card = slide6.shapes.add_shape(
    1, Inches(6.8), Inches(3.8), Inches(5.8), Inches(2.8))
pl2_card.fill.solid()
pl2_card.fill.fore_color.rgb = BG_MID
pl2_card.line.color.rgb = GREEN_SOFT
pl2_card.line.width = Pt(2)

add_text_box(slide6, 7.0, 3.9, 5.4, 0.4,
             "Plan 2:  Start with Easy-to-Evaluate Domains",
             font_size=18, color=GREEN_SOFT, bold=True, font_name="Calibri")

tf_pl2 = add_text_box(slide6, 7.0, 4.35, 5.4, 0.3,
                      "", font_size=14, color=WHITE, font_name="Calibri")
add_bullet(tf_pl2, "Data structures & algorithms have clear specs",
           font_size=14, color=WHITE)
add_bullet(tf_pl2, "Well-studied abstractions (AutoCLRS benchmark)",
           font_size=14, color=WHITE)
add_bullet(tf_pl2, "risemsr.github.io/blog/2026-03-06-autoclrs/",
           level=1, font_size=13, color=ACCENT2)
add_bullet(tf_pl2, "Measurable: proof lines, solver time, reuse count",
           font_size=14, color=WHITE)
add_bullet(tf_pl2, "Then transfer learned patterns to OS-level (Nanvix)",
           font_size=14, color=WHITE)

# Bottom tagline
add_text_box(slide6, 0.6, 6.7, 12, 0.4,
             "Strategy: prove the approach on algorithms first, then scale to systems code.",
             font_size=16, color=ACCENT, font_name="Calibri", alignment=PP_ALIGN.CENTER)


# Save
out_path = os.path.join(os.path.dirname(__file__), "set_theoretic_analysis.pptx")
prs.save(out_path)
print(f"Saved to {out_path}")
